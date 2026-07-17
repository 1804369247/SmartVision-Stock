# -*- coding: utf-8 -*-
"""根据《答辩要点.md》生成 SmartVision Stock 答辩 PPT 框架 (python-pptx)。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 配色 ----
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)   # 深蓝（主色）
BLUE   = RGBColor(0x2E, 0x6F, 0xB5)   # 蓝
LIGHT  = RGBColor(0xF2, 0xF6, 0xFB)   # 浅底
GREY   = RGBColor(0x55, 0x5F, 0x6B)   # 灰字
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0x7A, 0x1E)   # 橙强调

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def textbox(slide, x, y, w, h, lines, size=18, bold=False, color=NAVY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    """lines: list of (text, level) or str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in lines:
        if isinstance(item, tuple):
            txt, lvl = item
        else:
            txt, lvl = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.level = lvl
        run = p.add_run(); run.text = txt
        run.font.size = Pt(size - lvl * 2)
        run.font.bold = bold and lvl == 0
        run.font.color.rgb = color if lvl == 0 else GREY
        run.font.name = font
        p.space_after = Pt(6)
    return tb

def header(slide, title, sub=None):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.8))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEE); r2.font.name = "Microsoft YaHei"

def bullets(slide, items, top=Inches(1.5), size=18, left=Inches(0.7)):
    textbox(slide, left, top, SW - Inches(1.4), SH - top - Inches(0.4),
            items, size=size, color=NAVY)

# ============ 1. 封面 ============
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.6), SW, Inches(0.08), ACCENT)
textbox(s, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.6),
        [("SmartVision Stock", 0), ("智能仓储管理系统 —— 毕业设计答辩", 0)],
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.9), SW - Inches(2), Inches(1.6),
        [("基于 Vue 3 + Spring Boot 的 3D 可视化仓储管理系统", 0),
         ("答辩人：________    学号：________    导师：________", 0)],
        size=18, color=RGBColor(0xCF, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# ============ 2. 项目概述 ============
s = add_slide(); header(s, "一、项目概述", "系统定位与技术栈")
rows = [
    ("系统类型", "B/S 架构仓储管理系统（WMS），带 3D 可视化"),
    ("前端", "Vue 3 + Vite + Pinia + Element Plus + Three.js"),
    ("后端", "Spring Boot 2.7 + Spring Security + JWT + JPA + WebSocket"),
    ("数据库", "MySQL 8（生产）/ H2（测试）"),
    ("中间件", "Redis（缓存 + JWT 黑名单）、MinIO（对象存储）"),
    ("部署", "Docker + docker-compose 一键部署（nginx 托管前端）"),
    ("核心端口", "前端 8081，后端 8080"),
]
top = Inches(1.5)
for i, (k, v) in enumerate(rows):
    y = top + Inches(0.62) * i
    rect(s, Inches(0.7), y, Inches(2.4), Inches(0.52), BLUE if i % 2 == 0 else NAVY)
    textbox(s, Inches(0.8), y + Inches(0.04), Inches(2.2), Inches(0.46),
            [(k, 0)], size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(3.2), y, Inches(9.4), Inches(0.52), LIGHT)
    textbox(s, Inches(3.35), y + Inches(0.04), Inches(9.1), Inches(0.46),
            [(v, 0)], size=14, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# ============ 3. 核心功能模块 ============
s = add_slide(); header(s, "二、核心功能模块", "七大模块覆盖仓储全流程")
bullets(s, [
    "基础数据管理：商品、供应商、仓库、库位（含 3D 坐标）",
    "入库管理：入库单创建 → 审核 → 上架，生成货物实例（批次）",
    "出库管理：出库单创建 → 库存预占 → 审核 → 确认扣减，支持拣货",
    "库存管理：实时库存、调拨、移库、盘点、退货",
    "3D 可视化：Three.js 渲染仓库/货架/库位，实时反映库存状态",
    "实时通知：WebSocket 推送库存变更、到期预警",
    "系统管理：用户/角色/权限（RBAC）、操作日志、文件上传",
], size=20)

# ============ 4. 系统架构 ============
s = add_slide(); header(s, "三、系统架构", "口述用分层架构")
arch = (
    "[浏览器 Vue3 SPA]\n"
    "   │  REST(JWT) / WebSocket(子协议带 token)\n"
    "   ▼\n"
    "[Nginx] ──反向代理──► [Spring Boot 后端]\n"
    "                          ├─ Controller 层（REST API）\n"
    "                          ├─ Service 层（业务：出入库/预占/并发锁）\n"
    "                          ├─ Repository 层（JPA）\n"
    "                          ├─ WebSocket（实时推送）\n"
    "                          ▼\n"
    "              [MySQL]   [Redis]   [MinIO]\n"
    "              业务数据   缓存/黑名单  文件对象"
)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = arch
r.font.size = Pt(18); r.font.name = "Consolas"; r.font.color.rgb = NAVY

# ============ 5. 亮点1 库存预占 ============
s = add_slide(); header(s, "四、技术亮点 ①", "库存预占（Reserved Stock）机制 —— 解决“超卖”痛点")
bullets(s, [
    ("问题", 0),
    ("多个出库单同时想出同一批货，若只在最终确认时扣库存，会出现“账面够、实际不够”的超卖", 1),
    ("方案", 0),
    ("出库单审核阶段就逻辑锁定库存（写 stock_reservation 记录）", 1),
    ("可用量 = 库存数量 − 活动预占量", 1),
    ("预占带 TTL，@Scheduled 定时任务自动释放过期预占，避免“占而不用”锁死库存", 1),
    ("价值", 0),
    ("把库存冲突从“确认时”提前到“审核时”发现，符合真实仓储业务逻辑", 1),
], size=18)

# ============ 6. 亮点2 并发安全 ============
s = add_slide(); header(s, "四、技术亮点 ②", "并发安全 —— 悲观锁防止超卖 / 丢失更新")
bullets(s, [
    ("问题", 0),
    ("一个出库单可能拆多条明细扣减同一货物实例；并发确认多个单据时，读-改-写存在丢失更新", 1),
    ("方案", 0),
    ("出库确认前，对所有涉及的货物实例 id 先 distinct + sorted（固定加锁顺序防死锁）", 1),
    ("再逐个 findByIdForUpdate（SELECT … FOR UPDATE 悲观写锁）加锁", 1),
    ("预检查与扣减复用同一批已加锁实体", 1),
    ("验证", 0),
    ("@SpringBootTest 并发集成测试：8 线程确认库存 20（各扣 2）→ 剩余永不为负；10 线程并发预占不超卖", 1),
    ("全量 104 个测试 0 失败；为何不用乐观锁：历史 @Version 为 NULL 引发 NPE 且冲突频繁", 1),
], size=16)

# ============ 7. 亮点3 安全设计 ============
s = add_slide(); header(s, "四、技术亮点 ③", "安全设计")
bullets(s, [
    "认证：JWT（HS384 签名），Access + Refresh 双令牌，前端拦截器自动刷新",
    "登出即失效：JWT 黑名单，单机内存实现，多实例用 Redis 分布式黑名单",
    "收紧 CORS/权限：Spring Security 仅对 auth/swagger 等放行",
    "WebSocket 握手做 JWT 校验，且 CORS 不再用 * 通配兜底（为空直接拒绝）",
    "密钥管理：JWT 密钥、数据库密码全部走环境变量占位（${…}），代码零硬编码；缺密钥 fail-fast",
], size=18)

# ============ 8. 亮点4 工程化 ============
s = add_slide(); header(s, "四、技术亮点 ④", "实时性与工程化")
bullets(s, [
    "WebSocket 实时推送：库存变更事件驱动推送到 3D 前端；广播移出事务边界避免长事务",
    "缓存：Redis @Cacheable 缓存商品等热点数据",
    "可部署性：Docker 一键起全栈；/instances 等大列表接口做了分页，避免整表加载",
    "可测试性：H2 内存库 + test profile，集成测试不依赖外部 MySQL/Redis/MinIO",
], size=19)

# ============ 9. 演示流程 ============
s = add_slide(); header(s, "五、演示流程建议", "Demo 顺序（建议 8-10 分钟）")
bullets(s, [
    "1. 登录（admin/admin123）→ 展示 JWT 登录、进入首页",
    "2. 3D 仓库场景：Three.js 渲染货架/库位，点击库位看库存",
    "3. 入库：建入库单 → 审核 → 上架，3D 场景库位状态实时更新",
    "4. 出库 + 预占：建两个出库单占同一批货，展示“可用量”随预占减少；确认一单后库存扣减",
    "5. 超卖演示（加分项）：并发确认两单，展示悲观锁保证不扣成负数（可讲测试代替现场压测）",
    "6. 实时通知：另开窗口展示库存变更 WebSocket 推送",
    "7. 系统管理：用户/角色/权限、操作日志",
], size=17)

# ============ 10. Q&A ============
s = add_slide(); header(s, "六、可能的问题 & 参考回答", "Q&A 预案")
qa = [
    ("Q1 怎么防止超卖？", "业务层库存预占提前锁定 + 数据层悲观写锁保证原子性 + 固定加锁顺序防死锁 + 并发测试验证。"),
    ("Q2 为什么用悲观锁而非乐观锁？", "预占/出库确认冲突频繁，乐观锁重试成本高；且历史 @Version 为 NULL 引发过 NPE。"),
    ("Q3 Redis 起什么作用？", "缓存热点数据减轻 DB 压力；JWT 登出黑名单分布式共享，保证多实例登出即时生效；无 Redis 自动回退内存。"),
    ("Q4 WebSocket 如何鉴权？", "握手阶段子协议/Header 传 JWT，HandshakeInterceptor 校验；CORS 白名单不用 *，为空直接拒绝。"),
    ("Q5 3D 怎么实现？性能如何？", "Three.js 按库位坐标生成场景对象，库存变化更新材质/标签；注意 dispose 释放避免内存泄漏。"),
    ("Q6 数据一致性怎么保证？", "核心方法 @Transactional；WebSocket 广播移出事务边界，避免推送成功但事务回滚。"),
    ("Q7 测试覆盖如何？", "后端 104 个测试（单元+集成）含并发专项，全部通过；前端有 API/缓存/路由单测。"),
]
y = Inches(1.45)
for i, (q, a) in enumerate(qa):
    if i > 0 and i % 4 == 0:
        pass
    yy = y + Inches(0.82) * i
    textbox(s, Inches(0.7), yy, Inches(12), Inches(0.8),
            [(q, 0), (a, 1)], size=14, color=NAVY)

# ============ 11. 数据速记 ============
s = add_slide(); header(s, "七、项目数据速记", "背下来")
bullets(s, [
    "代码规模：本轮改造提交 254 个文件、约 9000+ 行新增",
    "测试：104 个后端测试，0 失败",
    "版本：v0.1.0（基座）→ v1.0.0（完整版），已打 tag 并推送 GitHub",
    "技术难点：库存预占、并发悲观锁、JWT 双令牌 + 分布式黑名单、3D 内存管理、Docker 一键部署",
], size=20)

# ============ 12. 可改进方向 ============
s = add_slide(); header(s, "八、可改进方向", "主动说，显得有反思")
bullets(s, [
    "多仓库调拨的库位直接关联可进一步完善（目前部分场景走间接关联）",
    "引入数据库迁移工具（Flyway/Liquibase）替代 ddl-auto:update，让 schema 变更可追溯",
    "高并发下可评估从悲观锁转向“预占 + 队列/分段库存”以提升吞吐",
    "增加更多端到端（E2E）自动化测试",
], size=20)

# ============ 13. 致谢 ============
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
textbox(s, Inches(1), Inches(2.8), SW - Inches(2), Inches(2),
        [("感谢各位老师批评指正！", 0), ("Q & A", 0)],
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

OUT = r"F:\SmartVision Stock\SmartVision_Stock_答辩PPT_框架.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
